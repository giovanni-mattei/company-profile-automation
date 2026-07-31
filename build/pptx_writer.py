from pptx import Presentation
from datetime import datetime
import os
from io import BytesIO

def update_text_fields(slide, text_fields):
    """Update text fields in a slide with provided values."""
    for shape in slide.shapes:
        if shape.name in text_fields and shape.has_text_frame:
            paragraph = shape.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].text = str(text_fields[shape.name])
            else:
                paragraph.add_run().text = str(text_fields[shape.name])

def replace_placeholders(slide, variables):
    """Replace placeholders in all text runs of a slide with corresponding values from a dictionary."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, value in variables.items():
                        if f"{{{placeholder}}}" in run.text:
                            run.text = run.text.replace(f"{{{placeholder}}}", str(value))


def generate_deck(company_data, analysis):
    
    template_path = "templates/deck.pptx"
    prs = Presentation(template_path)

    company_info = analysis["company_info"]
    company_summary = analysis["company_summary"]
    price_metrics = analysis["price_metrics"]
    price_chart = analysis["price_chart"]
    financial_highlights = analysis["financial_highlights"]
    financial_facts = analysis["financial_facts"]
    news = analysis["news"]

    current_date = datetime.now().strftime("%d %B %Y")
    company_name = company_info["company_name"]
    symbol = company_info["company_ticker"]

    # --- Slide 0: Title ---
    slide = prs.slides[0]
    replace_placeholders(slide, {"company_name": company_name, "current_date": current_date})
    
    # --- Slide 1: Company Overview ---
    slide = prs.slides[1]
    text_fields = {
        "name": str(company_info["company_name"]),
        "symbol": str(company_info["company_ticker"]),
        "sector": str(company_info["company_sector"]),
        "industry": str(company_info["company_industry"]),
        "market_cap": str(company_info["market_cap"]),
        "share_earnings": str(company_info["share_earnings"]),
        "share_price": str(company_info["share_price"]),
        "analyst_rating": str(company_info["analyst_rating"]),
        "beta": str(company_info["beta"]),
        "description": str(company_summary),
        "ltm_revenue": str(company_info["ltm_revenue"]),
        "revenue_growth": str(company_info["revenue_growth"]),
        "ltm_ebitda": str(company_info["ltm_ebitda"]),
        "ebitda_margin": str(company_info["ebitda_margin"])
    }
    update_text_fields(slide, text_fields)
    replace_placeholders(slide, {"company_name": company_name, "current_date": current_date})

    # --- Slide 2: Price Metrics ---
    slide = prs.slides[2]
    text_fields = {
        "current_price": f"{price_metrics['current_price']:.0f}",
        "mean_price": f"{price_metrics['mean_price']:.0f}",
        "max_price": f"{price_metrics['max_price']:.0f}",
        "75p_price": f"{price_metrics['75p_price']:.0f}",
        "50p_price": f"{price_metrics['50p_price']:.0f}",
        "25p_price": f"{price_metrics['25p_price']:.0f}",
        "min_price": f"{price_metrics['min_price']:.0f}",
        "std_price": f"{price_metrics['std_price']:.0f}",
        "analyst_expectation": f"{price_metrics['analyst_price']:.0f}"
    }
    update_text_fields(slide, text_fields)

    # Replace the price chart
    for shape in slide.shapes:
        if shape.name == "price_chart":
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(price_chart, left, top, width=width, height=height)
            break

    replace_placeholders(slide, {"company_name": company_name, "current_date": current_date})
    
    # --- Slide 3: Financial Highlights ---
    slide = prs.slides[3]
    
    years = list(financial_highlights["Revenue"].keys())[:3]
    
    text_fields = {}
    
    text_fields["year_1"] = str(years[0])
    text_fields["year_2"] = str(years[1])
    text_fields["year_3"] = str(years[2])
    
    text_fields["revenue_1"] = str(financial_highlights["Revenue"][years[0]])
    text_fields["revenue_2"] = str(financial_highlights["Revenue"][years[1]])
    text_fields["revenue_3"] = str(financial_highlights["Revenue"][years[2]])
    text_fields["revenue_m"] = financial_highlights["Revenue"].get("Avg growth", "")
    
    text_fields["gross_profit_1"] = str(financial_highlights["Gross Profit"][years[0]])
    text_fields["gross_profit_2"] = str(financial_highlights["Gross Profit"][years[1]])
    text_fields["gross_profit_3"] = str(financial_highlights["Gross Profit"][years[2]])
    text_fields["gross_profit_m"]     = financial_highlights["Gross Profit"].get("Avg margin", "")
    
    text_fields["operating_profit_1"] = str(financial_highlights["Operating Profit"][years[0]])
    text_fields["operating_profit_2"] = str(financial_highlights["Operating Profit"][years[1]])
    text_fields["operating_profit_3"] = str(financial_highlights["Operating Profit"][years[2]])
    text_fields["operating_profit_m"] = financial_highlights["Operating Profit"].get("Avg margin", "")
    
    text_fields["net_income_1"] = str(financial_highlights["Net Income"][years[0]])
    text_fields["net_income_2"] = str(financial_highlights["Net Income"][years[1]])
    text_fields["net_income_3"] = str(financial_highlights["Net Income"][years[2]])
    text_fields["net_income_m"] = financial_highlights["Net Income"].get("Avg margin", "")

    text_fields["invested_capital_1"] = str(financial_highlights["Invested Capital"][years[0]])
    text_fields["invested_capital_2"] = str(financial_highlights["Invested Capital"][years[1]])
    text_fields["invested_capital_3"] = str(financial_highlights["Invested Capital"][years[2]])
    text_fields["invested_capital_m"] = financial_highlights["Invested Capital"].get("Avg return", "")
    
    text_fields["debt_1"] = str(financial_highlights["Debt"][years[0]])
    text_fields["debt_2"] = str(financial_highlights["Debt"][years[1]])
    text_fields["debt_3"] = str(financial_highlights["Debt"][years[2]])
    text_fields["debt_m"] = financial_highlights["Debt"].get("Avg debt to capital ratio", "")

    text_fields["operating_cash_flow_1"] = str(financial_highlights["Operating Cash Flow"][years[0]])
    text_fields["operating_cash_flow_2"] = str(financial_highlights["Operating Cash Flow"][years[1]])
    text_fields["operating_cash_flow_3"] = str(financial_highlights["Operating Cash Flow"][years[2]])
    text_fields["operating_cash_flow_m"] = financial_highlights["Operating Cash Flow"].get("% of revenue", "")
    
    text_fields["investing_cash_flow_1"] = str(financial_highlights["Investing Cash Flow"][years[0]])
    text_fields["investing_cash_flow_2"] = str(financial_highlights["Investing Cash Flow"][years[1]])
    text_fields["investing_cash_flow_3"] = str(financial_highlights["Investing Cash Flow"][years[2]])
    text_fields["investing_cash_flow_m"] = financial_highlights["Investing Cash Flow"].get("% of revenue", "")
    
    text_fields["free_cash_flow_1"] = str(financial_highlights["Free Cash Flow"][years[0]])
    text_fields["free_cash_flow_2"] = str(financial_highlights["Free Cash Flow"][years[1]])
    text_fields["free_cash_flow_3"] = str(financial_highlights["Free Cash Flow"][years[2]])
    text_fields["free_cash_flow_m"] = financial_highlights["Free Cash Flow"].get("% growth", "")
            
    text_fields["comment_1"] = financial_facts[0]
    text_fields["comment_2"] = financial_facts[1]
    text_fields["comment_3"] = financial_facts[2]
    
    update_text_fields(slide, text_fields)
    replace_placeholders(slide, {"company_name": company_name, "current_date": current_date})

    # --- Slide 4: News ---
    slide = prs.slides[4]
    text_fields = {
        "news_name_1": news[0]['title'],
        "news_1": news[0]['summary'],
        "news_name_2": news[1]['title'],
        "news_2": news[1]['summary'],
        "news_name_3": news[2]['title'],
        "news_3": news[2]['summary']
    }
    update_text_fields(slide, text_fields)
    replace_placeholders(slide, {"company_name": company_name, "current_date": current_date})
    
    # --- Save the Presentation ---
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    current_date_filename = datetime.now().strftime("%d-%m-%Y")
    filename = f"{symbol}_{current_date_filename}.pptx"
    return {"file_content": buffer, "filename": filename}